import os
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt

def convert_pdf_to_pptx_images(pdf_path, pptx_path):
    print(f"Opening {pdf_path}...")
    try:
        doc = fitz.open(pdf_path)
        prs = Presentation()
        
        # Set slide dimensions to match PDF first page
        if len(doc) > 0:
            page = doc[0]
            # PDF points (1/72 inch) vs PPTX EMU (914400 per inch)
            # PyMuPDF rect width is in points. 1 point = 12700 EMUs
            prs.slide_width = int(page.rect.width * 12700)
            prs.slide_height = int(page.rect.height * 12700)
            print(f"Set presentation size to {page.rect.width}x{page.rect.height} points")

        # Blank layout
        blank_layout = prs.slide_layouts[6] 
        
        for i, page in enumerate(doc):
            print(f"Processing page {i+1}/{len(doc)}...")
            
            # Render page to image (zoom=2 for better quality)
            # matrix = fitz.Matrix(2, 2)
            pix = page.get_pixmap(dpi=150) # 150 DPI is usually good for screen
            
            img_path = f"temp_slide_{i}.png"
            pix.save(img_path)
            
            # Add slide
            slide = prs.slides.add_slide(blank_layout)
            
            # Add image covering whole slide
            slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
            
            # Cleanup temp image
            os.remove(img_path)
            
        prs.save(pptx_path)
        print(f"Successfully created {pptx_path}")
        
    except Exception as e:
        print(f"Error converting PDF: {e}")
        import traceback
        traceback.print_exc()

if __name__ == "__main__":
    pdf_file = "Nervestack Programming Language (Text Only).pdf"
    pptx_file = "Nervestack_Programming_Language_Visual.pptx"
    
    if os.path.exists(pdf_file):
        convert_pdf_to_pptx_images(pdf_file, pptx_file)
    else:
        print(f"File not found: {pdf_file}")
