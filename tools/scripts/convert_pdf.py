import os
from pypdf import PdfReader
from pptx import Presentation
from pptx.util import Inches, Pt

def convert_pdf_to_pptx(pdf_path, pptx_path):
    print(f"Reading {pdf_path}...")
    try:
        reader = PdfReader(pdf_path)
        prs = Presentation()
        
        # Calculate approximate chars per slide to avoid overflow
        # A standard slide might fit ~1000 chars comfortably at 12pt?
        # Let's just create one slide per page for now and use scrolling logic or small font if needed.
        # Actually, let's just dump page content to slide.
        
        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            
            # Create a blank slide
            slide_layout = prs.slide_layouts[1] # Title and Content
            slide = prs.slides.add_slide(slide_layout)
            
            # Title
            title = slide.shapes.title
            title.text = f"Page {i+1}"
            
            # Content
            content = slide.placeholders[1]
            content.text = text
            
            # Adjust font size to fit if possible
            for paragraph in content.text_frame.paragraphs:
                paragraph.font.size = Pt(14)  # Smaller font for full page text
                
        prs.save(pptx_path)
        print(f"Successfully created {pptx_path}")
        
    except Exception as e:
        print(f"Error converting PDF: {e}")

if __name__ == "__main__":
    pdf_file = "Nervestack Programming Language (Text Only).pdf"
    pptx_file = "Nervestack_Programming_Language.pptx"
    
    if os.path.exists(pdf_file):
        convert_pdf_to_pptx(pdf_file, pptx_file)
    else:
        print(f"File not found: {pdf_file}")
