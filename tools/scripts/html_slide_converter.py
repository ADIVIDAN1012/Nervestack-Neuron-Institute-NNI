
import asyncio
import os
from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Emu
from docx import Document
from docx.shared import Mm

# Configuration
HTML_FILE = "The_Story_of_Nervestack.html"
IMAGE_DIR = "slide_images"
OUTPUT_DOCX = "The_Story_of_Nervestack.docx"
OUTPUT_PPTX = "The_Story_of_Nervestack.pptx"

async def html_slides_to_images(html_path, output_dir):
    print(f"--- Extracting slides from {html_path} ---")
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
        
    image_files = []
    
    async with async_playwright() as p:
        browser = await p.chromium.launch()
        page = await browser.new_page()
        
        # Set viewport large enough to ensure responsive layout triggers desktop view
        await page.set_viewport_size({"width": 1920, "height": 3000}) # Height needs to be enough to render content? 
        # Actually for element screenshot, viewport width matters most for layout.
        
        abs_path = os.path.abspath(html_path)
        url = f"file:///{abs_path.replace(os.sep, '/')}"
        print(f"Loading {url}...")
        await page.goto(url)
        
        # Select all slide containers
        slides = await page.locator(".slide-container").all()
        print(f"Found {len(slides)} slides.")
        
        for i, slide in enumerate(slides):
            image_filename = f"slide_{i+1:02d}.png"
            image_path = os.path.join(output_dir, image_filename)
            
            # Screenshot the element directly
            await slide.screenshot(path=image_path)
            image_files.append(image_path)
            print(f"Saved {image_path}")
            
        await browser.close()
    return image_files

def images_to_docx(image_files, output_docx):
    print(f"--- Images to DOCX: {output_docx} ---")
    doc = Document()
    
    # Configure section to minimize margins
    section = doc.sections[0]
    section.page_height = Mm(297)
    section.page_width = Mm(210)
    section.left_margin = Mm(10)
    section.right_margin = Mm(10)
    section.top_margin = Mm(10)
    section.bottom_margin = Mm(10)
    
    for img_path in image_files:
        # Add image fitting the width of the page approximately (leaving room for margins)
        doc.add_picture(img_path, width=Mm(190))
        # Add a paragraph break
        doc.add_paragraph()
    
    doc.save(output_docx)
    print("DOCX saved.")

def images_to_pptx(image_files, output_pptx):
    print(f"--- Images to PPTX: {output_pptx} ---")
    prs = Presentation()
    
    # Set to 16:9 for modern layout
    prs.slide_width = Emu(914400 * 16)  # 16 inches
    prs.slide_height = Emu(914400 * 9)   # 9 inches
    
    for img_path in image_files:
        blank_slide_layout = prs.slide_layouts[6] 
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Fit image to slide
        # We assume the image is roughly 16:9. We can stretch to fit or fit within.
        # Given the CSS max-width 1280 and min-height 720, it's 16:9.
        slide.shapes.add_picture(
            img_path, 
            0, 
            0, 
            width=prs.slide_width, 
            height=prs.slide_height
        )
        
    prs.save(output_pptx)
    print("PPTX saved.")

async def main():
    if not os.path.exists(HTML_FILE):
        print(f"Error: {HTML_FILE} not found.")
        return

    # 1. HTML -> Images
    image_files = await html_slides_to_images(HTML_FILE, IMAGE_DIR)
    
    if not image_files:
        print("No images generated. Exiting.")
        return

    # 2. Images -> DOCX
    images_to_docx(image_files, OUTPUT_DOCX)
    
    # 3. Images -> PPTX
    images_to_pptx(image_files, OUTPUT_PPTX)
    
    print("All conversions complete.")

if __name__ == "__main__":
    asyncio.run(main())
